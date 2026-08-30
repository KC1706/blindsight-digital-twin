"""Build the Detailed Business Proposal as a .pptx (Blindsight / Team DarkShield).

Mirrors docs/BUSINESS_PROPOSAL.md. Dark theme, Accenture-purple accent.
Usage: python scripts/build_proposal_pptx.py submission/Blindsight_Business_Proposal.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

BG = RGBColor(0x0B, 0x0B, 0x10)
PANEL = RGBColor(0x1B, 0x1B, 0x22)
TEXT = RGBColor(0xF2, 0xF2, 0xF6)
DIM = RGBColor(0x9A, 0x9A, 0xA6)
PUR = RGBColor(0xA1, 0x00, 0xFF)
PUR2 = RGBColor(0xC4, 0x6B, 0xFF)
GOOD = RGBColor(0x7E, 0xE7, 0x87)
WARN = RGBColor(0xE3, 0xB3, 0x41)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         bullet=False, space=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    run = p.add_run()
    run.text = ("• " + text) if bullet else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return p


def kicker(s, text):
    tf = box(s, 0.7, 0.5, 12, 0.5)
    para(tf, text.upper(), 12, PUR2, bold=True, first=True)


def title(s, text):
    tf = box(s, 0.7, 0.9, 12, 1.1)
    para(tf, text, 30, TEXT, bold=True, first=True)


def card(s, l, t, w, h, border=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                             Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = PANEL
    shp.line.color.rgb = PUR if border else PANEL
    shp.line.width = Pt(1.5 if border else 0.5)
    shp.shadow.inherit = False
    return shp


def statcard(s, l, t, w, stat, cap, color=GOOD):
    card(s, l, t, w, 1.7, border=True)
    tf = box(s, l + 0.25, t + 0.2, w - 0.5, 1.3)
    para(tf, stat, 30, color, bold=True, first=True, space=2)
    para(tf, cap, 11, DIM)


def table(s, l, t, w, rows, colw, header=True, fs=12):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w),
                            Inches(0.4 * nr)).table
    for ci, cw in enumerate(colw):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2a, 0x18, 0x3f) if (header and ri == 0) else PANEL
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            cell.text_frame.word_wrap = True
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(fs)
            r.font.name = "Segoe UI"
            r.font.bold = (header and ri == 0)
            r.font.color.rgb = PUR2 if (header and ri == 0) else TEXT


def footer(s, txt):
    tf = box(s, 0.7, 7.0, 12, 0.35)
    para(tf, txt, 9, DIM, first=True)


# ---- 1 · title ----
s = slide()
card(s, 0.7, 0.7, 0.9, 0.9, border=True)
tf = box(s, 0.7, 0.62, 0.9, 0.9, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "B", 34, TEXT, bold=True, first=True, align=PP_ALIGN.CENTER)
tf = box(s, 0.7, 2.1, 12, 3)
para(tf, "TEAM DARKSHIELD · AIC 2026 · TRACK 4 · DIGITALTWIN.AI", 12, PUR2, bold=True, first=True)
para(tf, "Blindsight", 46, TEXT, bold=True, space=4)
para(tf, "Detailed Business Proposal", 22, PUR2, bold=True, space=10)
para(tf, "A digital twin that sees the stations it can't measure —\n"
        "the line is the sensor.", 16, DIM)
footer(s, "Round 2 · working prototype · retrofit-free, read-only")

# ---- 2 · problem framing ----
s = slide(); kicker(s, "1 · Problem framing"); title(s, "The line is half blind — where it matters most")
tf = box(s, 0.7, 2.0, 12, 3.4)
para(tf, "Assembly lines mix legacy and modern equipment, so sensor coverage is uneven. "
        "Manual stations — often the true constraint — are dark: ~25% of the prototype's "
        "40-station line (10 stations), up to ~35% on a representative line.", 14, TEXT, first=True)
para(tf, "MES / SCADA dashboards show only what's wired, and only after the stop — a rear-view mirror.", 13, TEXT, bullet=True)
para(tf, "APM / predictive-maintenance tools need a sensor on the asset — no help for the dark stations.", 13, TEXT, bullet=True)
para(tf, "Digital-twin / simulation platforms are offline, design-time — not live inference from partial data.", 13, TEXT, bullet=True)
para(tf, "Retrofitting means scarce maintenance windows and capex.", 13, TEXT, bullet=True)
statcard(s, 0.7, 5.2, 4.6, "≈ ₹140 Cr", "cost of line-blindness / yr (base; ₹29–549 Cr)", WARN)
footer(s, "docs/COST_OF_BLINDNESS.md · MARKET_LANDSCAPE.md")

# ---- 3 · solution ----
s = slide(); kicker(s, "2 · Solution design"); title(s, "\"The line is the sensor\" — five moves")
tf = box(s, 0.7, 2.0, 12, 4.4)
para(tf, "Turn the sensors a plant already has into inference for the ones it doesn't, then close the loop:", 14, TEXT, first=True, space=10)
for n, t in [
    ("Virtual sensing", "blocked/starved propagation + travel-time tomography estimate every station's cycle time — dark or not — with error bars."),
    ("Probabilistic forecast", "2-hour Monte-Carlo: P(binding constraint), ETA, expected units lost, with bands."),
    ("Prescriptive action", "counterfactual re-simulation ranks levers by expected units recovered."),
    ("Defect containment", "SPC signal + scan timestamps name the exact VINs in the drift window, before audit."),
    ("Trust & value-of-information", "every forecast backtested; ranks which dark stations to instrument first."),
]:
    p = tf.add_paragraph(); p.space_after = Pt(7)
    r = p.add_run(); r.text = n + " — "; r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = PUR2; r.font.name = "Segoe UI"
    r2 = p.add_run(); r2.text = t; r2.font.size = Pt(13); r2.font.color.rgb = TEXT; r2.font.name = "Segoe UI"
footer(s, "Retrofit-free and read-only — no PLC writes, no new hardware · PLAN.md, docs/METHODS.md")

# ---- 4 · target users ----
s = slide(); kicker(s, "3 · Target users"); title(s, "Three stakeholders, three views — already built")
table(s, 0.7, 2.2, 12, [
    ["Persona", "Job-to-be-done", "View"],
    ["Floor supervisor", "Know the real constraint right now, without walking the line", "Live map + one recommended action"],
    ["Plant manager", "See the 2-hour forecast and a ranked, costed action", "Forecast bands, ranked levers, per-variant work"],
    ["Leadership", "Decide whether to fund a pilot/scale with an ROI number", "ROI tile + VoI-ranked instrumentation roadmap"],
], [2.6, 5.4, 4.0], fs=13)
footer(s, "docs/PERSONAS.md · all three views implemented in web/")

# ---- 5 · business case ----
s = slide(); kicker(s, "4 · Business case & impact"); title(s, "Retrofit-free value, in the first month")
statcard(s, 0.7, 2.1, 3.9, "~₹64 Cr", "addressable value / yr · one line (base, ramped)", GOOD)
statcard(s, 4.75, 2.1, 3.9, "< 1 month", "payback (base) — implementation cost alone", PUR2)
statcard(s, 8.8, 2.1, 3.85, "₹0", "new hardware — read-only tap on scans + PLC", GOOD)
tf = box(s, 0.7, 4.0, 12, 0.6)
para(tf, "Value = capture_rate × cost-of-blindness + deferred capex.  Priced per line: ₹8.5 lakh/mo + ~₹34 lakh one-time.", 13, DIM, first=True)
table(s, 0.7, 4.7, 12, [
    ["Scenario", "Annual value (ramped)", "Year-1 cost", "Payback"],
    ["Low — the number to lead with", "₹7.2 Cr", "₹1.03 Cr", "~5 months"],
    ["Base", "₹64 Cr", "₹1.36 Cr", "< 1 month"],
    ["High", "₹332 Cr", "₹2.04 Cr", "< 1 month"],
], [4.2, 3.2, 2.3, 2.3], fs=12)
footer(s, "docs/ROI_MODEL.md · PRICING.md · PILOT_KPIS.md")

# ---- 6 · roadmap ----
s = slide(); kicker(s, "5 · Phased roadmap"); title(s, "Prove it, pilot it, scale it")
table(s, 0.7, 2.3, 12, [
    ["Phase", "Scope", "Exit gate"],
    ["1 — PoC (this prototype)", "One simulated line, mechanism proven end-to-end", "Validation targets met; live demo runs clean"],
    ["2 — Pilot", "One real line, read-only tap, shadow-mode before go-live", "KPI thresholds met on real data; reference case"],
    ["3 — Scale", "Multi-line, multi-site; VoI-driven instrumentation", "Steady state, tracked per-line via the KPI sheet"],
], [3.0, 5.0, 4.0], fs=12.5)
tf = box(s, 0.7, 5.4, 12, 0.8)
para(tf, "No license fee during Pilot; per-line SaaS tiers at Scale. Value-of-information turns the twin's own "
        "uncertainty into a costed plan: instrument the 3 highest-impact dark stations first, not all of them.", 13, DIM, first=True)
footer(s, "docs/GTM_ROADMAP.md · PRICING.md")

# ---- 7 · risks ----
s = slide(); kicker(s, "6 · Key risks & mitigations"); title(s, "The risks a skeptic will raise — and the answers")
table(s, 0.7, 2.2, 12, [
    ["Risk", "Mitigation"],
    ["Dark-station estimates too uncertain to act on", "Error bars everywhere; abstain below a confidence threshold; VoI says where a cheap sensor removes the doubt"],
    ["False alarms erode floor trust", "Backtested false-alarm rate (0% on clean baselines); conservative thresholds; shadow-mode first"],
    ["OT / security concerns", "Read-only by construction — no PLC write path; DMZ-bridged tap, IEC-62443-aligned segmentation"],
    ["Site-to-site variation", "Per-site recalibration before go-live; hierarchical priors transfer structure, not exact values"],
    ["Change management (staff ignore it)", "One clear action, not a dashboard to interpret; overrides captured as feedback"],
], [3.7, 8.3], fs=11.5)
footer(s, "docs/RISK_REGISTER.md — full register + compliance posture")

# ---- 8 · validation ----
s = slide(); kicker(s, "7 · Validation — why the numbers can be trusted"); title(s, "The twin scoring itself")
statcard(s, 0.7, 2.2, 3.0, "3.3%", "dark-station cycle-time MAPE (no sensor)", GOOD)
statcard(s, 3.85, 2.2, 3.0, "0%", "false-alarm rate · 12 clean shifts", GOOD)
statcard(s, 7.0, 2.2, 3.0, "~10×", "forecast CRPS vs persistence baseline", GOOD)
statcard(s, 10.15, 2.2, 2.5, "100%", "error-bar coverage (80% band)", GOOD)
tf = box(s, 0.7, 4.3, 12, 1.5)
para(tf, "Every predictive claim is backtested, not asserted — split honestly by instrumented vs dark stations, "
        "on held-out shifts. On takt_slip_s14, the dark bottleneck S14 is localized with no sensor on it and its "
        "cycle estimated at 68.5s (true 67.4s), with ~9 minutes of median lead time before the first unit is lost.", 14, TEXT, first=True)
para(tf, "Reproducible:  python -m engine.validate", 12, DIM)
footer(s, "docs/VALIDATION.md · 27 automated tests, CI green")

# ---- 9 · close ----
s = slide()
card(s, 5.7, 1.6, 0.95, 0.95, border=True)
tf = box(s, 5.7, 1.52, 0.95, 0.95, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "B", 34, TEXT, bold=True, first=True, align=PP_ALIGN.CENTER)
tf = box(s, 1, 3.0, 11.33, 2.5, anchor=MSO_ANCHOR.TOP)
para(tf, "No retrofit. No new sensors.\nJust the line, finally telling you what it knows.", 26, TEXT, bold=True, first=True, align=PP_ALIGN.CENTER)
para(tf, "Team DarkShield · Blindsight · the line is the sensor", 14, DIM, align=PP_ALIGN.CENTER)

out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("submission/Blindsight_Business_Proposal.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"wrote {out} ({out.stat().st_size // 1024} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
